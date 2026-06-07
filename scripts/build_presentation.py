#!/usr/bin/env python3
"""
build_presentation.py — Gera PPTX de apresentação do projeto de fábrica.

Entrada:  data/projeto.json + data/resultados_calculo.json
Saída:    01_apresentacao/apresentacao_tramontina_22399036.pptx
"""

import json
import tempfile
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT    = Path(__file__).parent.parent
DATA    = ROOT / "data"
RENDERS = ROOT / "06_dashboard" / "renders"
ASSETS  = ROOT / "04_fontes" / "assets_tramontina"
OUT_DIR = ROOT / "01_apresentacao"

# ── palette ──────────────────────────────────────────────────────────────────
BLUE  = RGBColor(0x31, 0x5f, 0x9f)
GREEN = RGBColor(0x2f, 0x7d, 0x57)
GOLD  = RGBColor(0xb0, 0x78, 0x20)
RED   = RGBColor(0xa6, 0x3d, 0x40)
INK   = RGBColor(0x17, 0x21, 0x2b)
MUTED = RGBColor(0x5d, 0x69, 0x77)
PAPER = RGBColor(0xf8, 0xfa, 0xfc)
WHITE = RGBColor(0xff, 0xff, 0xff)
LINE  = RGBColor(0xd7, 0xdd, 0xe5)
PURP  = RGBColor(0x6a, 0x4c, 0x93)
LGRAY = RGBColor(0xed, 0xf2, 0xf7)

TYPE_COLOR = {
    "operacao":    GREEN,
    "transporte":  BLUE,
    "inspecao":    GOLD,
    "armazenagem": PURP,
    "espera":      RED,
}
TYPE_LABEL = {
    "operacao":    "Operacao",
    "transporte":  "Transporte",
    "inspecao":    "Inspecao",
    "armazenagem": "Armazenagem",
    "espera":      "Espera",
}

# ── geometry (inches) ────────────────────────────────────────────────────────
W  = 10.0
H  = 5.625
MX = 0.38
HH = 0.60
FY = H - 0.18

# ── low-level primitives ─────────────────────────────────────────────────────

def _rect(slide, x, y, w, h, fill=None, line_color=None):
    sp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    f = sp.fill
    if fill:
        f.solid(); f.fore_color.rgb = fill
    else:
        f.background()
    ln = sp.line
    if line_color:
        ln.color.rgb = line_color; ln.width = Pt(0.5)
    else:
        ln.fill.background()
    return sp


def _txt(slide, text, x, y, w, h, size=12, bold=False, italic=False,
         color=INK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _lines(slide, items, x, y, w, h, size=10, color=INK, spacing_pt=2, bold=False):
    """Multi-paragraph textbox — one entry per paragraph."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if spacing_pt:
            p.space_before = Pt(spacing_pt)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def _img(slide, path, x, y, w, h):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def _header(slide, num, title, color=BLUE):
    _rect(slide, 0, 0, W, HH, fill=color)
    _txt(slide, f"{num:02d}", MX, 0.04, 0.42, HH - 0.08,
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, MX + 0.46, 0.12, 0.02, HH - 0.24, fill=WHITE)
    _txt(slide, title, MX + 0.56, 0.09, W - MX - 0.7, HH - 0.18,
         size=17, bold=True, color=WHITE)


def _footer(slide):
    _txt(slide, "Projeto de Fabrica · Tramontina Kit Churrasco 22399036 · Grupo A",
         MX, FY, W - 2 * MX, 0.18, size=7, color=MUTED, align=PP_ALIGN.CENTER)


def _kpi(slide, label, value, unit, x, y, w, h, color=BLUE):
    _rect(slide, x, y, w, h, fill=WHITE, line_color=LINE)
    _rect(slide, x, y, 0.05, h, fill=color)
    _txt(slide, value,
         x + 0.1, y + 0.04, w - 0.14, h * 0.52,
         size=min(22, int(w * 18)), bold=True, color=color, align=PP_ALIGN.CENTER)
    _txt(slide, unit,
         x + 0.1, y + h * 0.52, w - 0.14, h * 0.25,
         size=7, color=MUTED, align=PP_ALIGN.CENTER)
    _txt(slide, label,
         x + 0.1, y + h * 0.75, w - 0.14, h * 0.23,
         size=8, bold=True, color=INK, align=PP_ALIGN.CENTER)


def _tbl(slide, headers, rows, x, y, w, h, col_w=None, row_size=8):
    n_rows, n_cols = len(rows) + 1, len(headers)
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h)
    ).table
    if col_w:
        for ci, cw in enumerate(col_w):
            tbl.columns[ci].width = Inches(cw)

    def _fmt(cell, text, bold=False, fg=INK, bg=None):
        cell.text = str(text)
        if bg:
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(row_size)
                run.font.bold = bold
                run.font.color.rgb = fg

    for ci, h_txt in enumerate(headers):
        _fmt(tbl.cell(0, ci), h_txt, bold=True, fg=WHITE, bg=BLUE)
    for ri, row in enumerate(rows):
        bg = PAPER if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            _fmt(tbl.cell(ri + 1, ci), val, bg=bg)
    return tbl


# ── SVG → PNG ────────────────────────────────────────────────────────────────

def _svg_to_png(svg: Path, png: Path):
    from playwright.sync_api import sync_playwright
    edge = Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe")
    kwargs = {"headless": True}
    if edge.exists():
        kwargs["executable_path"] = str(edge)
    with sync_playwright() as p:
        browser = p.chromium.launch(**kwargs)
        try:
            page = browser.new_page(viewport={"width": 1800, "height": 1000})
            page.goto("file:///" + str(svg).replace("\\", "/"), wait_until="load")
            page.screenshot(path=str(png))
        finally:
            browser.close()


def _render_pngs(tmp: Path) -> dict:
    out = {}
    for name in ("fluxograma_render", "layout_render", "mapofluxograma_render"):
        svg = RENDERS / f"{name}.svg"
        png = tmp / f"{name}.png"
        print(f"  {name}.svg -> PNG ...")
        _svg_to_png(svg, png)
        out[name] = png
    return out


# ── blank slide helper ───────────────────────────────────────────────────────

def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ── slide builders ────────────────────────────────────────────────────────────

def s01_cover(prs, proj, res):
    slide = _slide(prs)
    members = proj["metadata"]["group_members"]

    _rect(slide, 0, 0, W, 0.10, fill=BLUE)

    _txt(slide, "PROJETO DE FABRICA", MX, 0.20, 5.8, 0.34,
         size=10, bold=True, color=MUTED)
    _txt(slide, "Kit para Churrasco Tramontina", MX, 0.50, 5.8, 0.55,
         size=24, bold=True, color=BLUE)
    _txt(slide, "com Laminas em Aco Inox e Cabos em Madeira Natural — 3 Pecas",
         MX, 1.00, 5.5, 0.50, size=14, color=INK)
    _rect(slide, MX, 1.58, 1.1, 0.03, fill=GOLD)
    _txt(slide, "SKU 22399036 · Tramontina · Arranjo Fisico Industrial — Grupo A",
         MX, 1.70, 5.5, 0.26, size=9, color=MUTED)

    _txt(slide, "Integrantes:", MX, 2.10, 3.5, 0.24, size=9, bold=True, color=INK)
    names = [m["name"] for m in members]
    _lines(slide, names, MX, 2.32, 3.5, 1.50, size=11, color=INK, spacing_pt=1)

    kpi_y = 4.12
    _kpi(slide, "Meta semanal", "1.000", "kits bons", MX,            kpi_y, 1.52, 0.88, BLUE)
    _kpi(slide, "Area fabrica",  "384 m2",   "24 x 16 m", MX + 1.62, kpi_y, 1.52, 0.88, GREEN)
    _kpi(slide, "Gargalo CNC",   "2x RTC.1313", "Router CNC", MX + 3.24, kpi_y, 1.52, 0.88, GOLD)

    _img(slide, ASSETS / "22399036_produto_principal_G.jpg", 6.4, 0.08, 3.25, 5.2)
    _footer(slide)


def s02_objetivos(prs, proj):
    slide = _slide(prs)
    _header(slide, 2, "Objetivos do Projeto")
    _footer(slide)

    objetivos = proj["project_objectives"]
    cy = HH + 0.14
    row_h = (FY - cy - 0.05) / len(objetivos)

    for i, obj in enumerate(objetivos):
        oy = cy + i * row_h
        bh = row_h - 0.06
        _rect(slide, MX, oy + 0.03, 0.36, bh, fill=BLUE)
        _txt(slide, str(i + 1), MX, oy + 0.03, 0.36, bh,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, obj, MX + 0.44, oy + 0.06, W - MX * 2 - 0.48, bh - 0.08,
             size=11.5, color=INK, wrap=True)


def s03_produto(prs, proj):
    slide = _slide(prs)
    _header(slide, 3, "Produto — Kit Churrasco Tramontina 22399036")
    _footer(slide)

    p = proj["product"]
    dims = p["package_dimensions_cm"]
    cy = HH + 0.14
    lw = 5.5

    _txt(slide, "Composicao do kit", MX, cy, lw, 0.24, size=9, bold=True, color=MUTED)
    _txt(slide, "Faca Chef 8\" · Garfo Trinchante · Tabua Retangular Macaranduba",
         MX, cy + 0.22, lw, 0.30, size=12, bold=True, color=INK)

    _txt(slide, "Materiais", MX, cy + 0.60, lw, 0.22, size=9, bold=True, color=MUTED)
    _lines(slide, ["• " + m for m in p["official_materials"]],
           MX, cy + 0.82, lw - 0.1, 0.90, size=9.5, color=INK, spacing_pt=1)

    _txt(slide, "Embalagem",  MX, cy + 1.82, lw, 0.22, size=9, bold=True, color=MUTED)
    _txt(slide, f"{dims['height']} x {dims['width']} x {dims['length']} cm  |  {p['package_weight_kg']} kg",
         MX, cy + 2.04, lw, 0.28, size=12, color=INK)

    sub = [
        ("22315008_faca_chef_8.jpg",      'Faca Chef 8"'),
        ("22330000_garfo_trinchante.jpg",  "Garfo Trinchante"),
        ("13102152_tabua_retangular.jpg",  "Tabua Macaranduba"),
    ]
    iw = (lw - 0.15) / 3
    iy = cy + 2.42
    for j, (fn, lbl) in enumerate(sub):
        ix = MX + j * (iw + 0.07)
        _img(slide, ASSETS / fn, ix, iy, iw - 0.05, 1.0)
        _txt(slide, lbl, ix, iy + 1.02, iw - 0.05, 0.2,
             size=7, color=MUTED, align=PP_ALIGN.CENTER)

    _img(slide, ASSETS / "22399036_item_aberto_G.jpg", 5.95, cy, 3.65, 3.90)


def s04_bom(prs, proj):
    slide = _slide(prs)
    _header(slide, 4, "Estrutura Pai-Filho — Arvore do Produto (BOM)")
    _footer(slide)

    bom = proj["bom"]
    fazer   = [b for b in bom if b["make_or_buy"] == "Fazer"]
    comprar = [b for b in bom if b["make_or_buy"] == "Comprar"]
    cy = HH + 0.14

    # Parent
    _rect(slide, MX, cy, W - 2 * MX, 0.50, fill=BLUE)
    _txt(slide, "Kit 22399036 — Kit para Churrasco Tramontina 3 Pecas (embalado para envio)",
         MX + 0.15, cy + 0.05, W - 2 * MX - 0.3, 0.40,
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(slide, "v", W / 2 - 0.1, cy + 0.50, 0.2, 0.26, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    child_y = cy + 0.78
    half_w  = (W - 2 * MX - 0.18) / 2

    def _group(items, gx, header_text, color, item_h):
        _rect(slide, gx, child_y, half_w, 0.28, fill=color)
        _txt(slide, header_text, gx + 0.08, child_y + 0.04, half_w - 0.16, 0.22,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for k, item in enumerate(items):
            iy = child_y + 0.30 + k * item_h
            _rect(slide, gx, iy, half_w, item_h - 0.04, fill=WHITE, line_color=color)
            _rect(slide, gx, iy, 0.05, item_h - 0.04, fill=color)
            _txt(slide, item["component"], gx + 0.12, iy + 0.04, half_w - 0.18, item_h - 0.12,
                 size=9, color=INK, wrap=True)

    available_h = FY - child_y - 0.35
    _group(fazer,   MX,                   f"FABRICAR INTERNAMENTE ({len(fazer)} itens)",
           GREEN, (available_h - 0.28) / len(fazer))
    _group(comprar, MX + half_w + 0.18,   f"COMPRAR DE FORNECEDOR ({len(comprar)} itens)",
           GOLD,  (available_h - 0.28) / len(comprar))

    _txt(slide, "Itens fabricados internamente para controle de qualidade e acabamento. "
                "Itens de embalagem adquiridos de fornecedores especializados.",
         MX, FY - 0.26, W - 2 * MX, 0.22, size=8, color=MUTED, italic=True)


def s05_mercado(prs, proj):
    slide = _slide(prs)
    _header(slide, 5, "Segmentos de Mercado")
    _footer(slide)

    segments = proj["market_segments"]
    colors = [BLUE, GREEN, GOLD, RED]
    cy = HH + 0.16
    cw = (W - 2 * MX - 0.18) / 2
    ch = (FY - cy - 0.14) / 2

    for i, seg in enumerate(segments):
        col, row = i % 2, i // 2
        sx = MX + col * (cw + 0.18)
        sy = cy + row * (ch + 0.14)
        c = colors[i]
        _rect(slide, sx, sy, cw, ch, fill=WHITE, line_color=c)
        _rect(slide, sx, sy, cw, 0.28, fill=c)
        _txt(slide, seg["name"], sx + 0.1, sy + 0.03, cw - 0.2, 0.24,
             size=10, bold=True, color=WHITE)
        _txt(slide, seg["justification"],
             sx + 0.1, sy + 0.32, cw - 0.2, ch - 0.38,
             size=8.5, color=INK, wrap=True)


def s06_meta(prs, proj, res):
    slide = _slide(prs)
    _header(slide, 6, "Meta Semanal de Producao")
    _footer(slide)

    pr     = proj["premises"]
    demand = res["demand"]
    cy = HH + 0.16

    kw = (W - 2 * MX - 0.44) / 5
    kh = 1.08
    kpis = [
        ("Meta boa",    f"1.000",  "kits/semana", GREEN),
        ("Dem. bruta",  str(demand["input_kits_per_week"]), "kits/semana", BLUE),
        ("Horas uteis", f"{demand['useful_hours_per_week']:.0f} h", "por semana", GOLD),
        ("Ritmo medio", f"{demand['required_average_rate_kits_per_hour']:.1f}", "kits/h", RED),
        ("Yield final", f"{int(pr['final_good_yield']*100)}%", "rendimento", PURP),
    ]
    for j, (lbl, val, unit, c) in enumerate(kpis):
        _kpi(slide, lbl, val, unit, MX + j * (kw + 0.11), cy, kw, kh, c)

    tbl_y = cy + kh + 0.18
    headers = ["Premissa", "Valor", "Fonte / Justificativa"]
    rows = [
        ["Meta semanal de kits bons",  "1.000 kits/semana",
         "Premissa do projeto"],
        ["Dias uteis / semana",         f"{pr['work_days_per_week']} dias",
         "Premissa"],
        ["Turnos / dia",                f"{pr['shifts_per_day']} turno",
         "Premissa"],
        ["Horas uteis / turno",
         f"{pr['useful_hours_per_shift']} h uteis de {pr['scheduled_hours_per_shift']} programadas",
         "Premissa — deduz pausa e manutencao"],
        ["Eficiencia geral",            f"{int(pr['general_efficiency']*100)}%",
         "Premissa de engenharia"],
        ["Confiabilidade equip.",       f"{int(pr['equipment_reliability']*100)}%",
         "Premissa de engenharia"],
        ["Rendimento final",            f"{int(pr['final_good_yield']*100)}%",
         "Premissa de engenharia — kits bons / producao bruta"],
    ]
    tbl_h = FY - tbl_y - 0.05
    _tbl(slide, headers, rows, MX, tbl_y, W - 2 * MX, tbl_h,
         col_w=[3.8, 2.4, 3.45], row_size=8)


def s07_tabela1(prs, proj):
    slide = _slide(prs)
    _header(slide, 7, "Tabela 1 — Componentes (BOM)")
    _footer(slide)

    cy = HH + 0.18
    headers = ["Componente", "Qtd.", "Un.", "Fazer / Comprar", "Justificativa"]
    rows = [
        [b["component"], b["quantity"], b["unit"],
         b["make_or_buy"], b["source_or_premise"]]
        for b in proj["bom"]
    ]
    _tbl(slide, headers, rows, MX, cy, W - 2 * MX, FY - cy - 0.05,
         col_w=[3.8, 0.42, 0.62, 1.28, 3.53], row_size=8)


def s08_processos(prs, proj):
    slide = _slide(prs)
    _header(slide, 8, "Processo Produtivo — Tabela 2 (26 processos)")
    _footer(slide)

    procs = proj["processes"]
    counts = Counter(p["type"] for p in procs)
    type_order = ["operacao", "inspecao", "transporte", "armazenagem", "espera"]
    cy = HH + 0.14

    chip_w = (W - 2 * MX - 0.40) / 5
    chip_h = 0.70
    for i, t in enumerate(type_order):
        cx = MX + i * (chip_w + 0.10)
        c = TYPE_COLOR[t]
        _rect(slide, cx, cy, chip_w, chip_h, fill=c)
        _txt(slide, str(counts.get(t, 0)), cx, cy + 0.02, chip_w, 0.36,
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, TYPE_LABEL[t], cx, cy + 0.38, chip_w, 0.28,
             size=8.5, color=WHITE, align=PP_ALIGN.CENTER)

    _txt(slide, "Total: 26 processos distribuidos em 5 setores de producao",
         MX, cy + chip_h + 0.08, W - 2 * MX, 0.22,
         size=9, bold=True, color=INK, align=PP_ALIGN.CENTER)

    tbl_y = cy + chip_h + 0.36
    headers = ["Setor", "Processos incluidos", "Recursos principais"]
    phases = [
        ("Recebimento (P1-P2)",    "Inspecao de entrada + Armazenagem de MP",
         "Doca, paleteira, balanca, estantes"),
        ("Setor Metal (P3-P10)",   "8 processos: corte laser, TT, polimento, afiacao",
         "Laser fibra CNC, Forno TT, Politriz, Afiador"),
        ("Setor Madeira (P11-P17)","7 processos: corte, usinagem CNC, lixamento, acabamento",
         "Esquadrejadeira, Router CNC (2x), Lixadeira"),
        ("Montagem (P18-P24)",     "7 processos: rebitagem, kit, blister, embalagem final",
         "Rebitadeira pneumatica, Bancada, Seladora Blister"),
        ("Expedicao (P25-P26)",    "Armazenagem de PA + Transporte para expedicao",
         "Porta-paletes, paleteira, doca de saida"),
    ]
    _tbl(slide, headers, phases, MX, tbl_y, W - 2 * MX, FY - tbl_y - 0.05,
         col_w=[2.0, 3.6, 4.05], row_size=8.5)


def s09_fluxograma(prs, renders):
    slide = _slide(prs)
    _header(slide, 9, "Fluxograma do Processo — 26 Etapas")
    _footer(slide)
    png = renders.get("fluxograma_render")
    if png and Path(png).exists():
        iy = HH + 0.08
        _img(slide, png, MX, iy, W - 2 * MX, FY - iy - 0.05)
    else:
        _txt(slide, "[Fluxograma nao disponivel]", MX, HH + 1.0, W - 2 * MX, 0.5,
             size=14, color=MUTED, align=PP_ALIGN.CENTER)


def s10_equipamentos(prs, res):
    slide = _slide(prs)
    _header(slide, 10, "Tabela 3 — Equipamentos Selecionados")
    _footer(slide)

    cy = HH + 0.16
    headers = ["Equipamento", "Fornecedor / Modelo", "Taxa efetiva (kits/h)", "Qtd.", "Utilizacao"]
    rows = [
        [e["type"],
         f"{e['supplier']} {e['model']}",
         f"{e['effective_rate_kits_per_hour']:.1f}",
         str(e["required_quantity"]),
         f"{e['utilization']*100:.1f}%"]
        for e in res["equipment_capacity"]
    ]
    _tbl(slide, headers, rows, MX, cy, W - 2 * MX, FY - cy - 0.05,
         col_w=[2.55, 2.8, 1.35, 0.48, 0.97], row_size=7.5)


def s11_calculo(prs, res):
    slide = _slide(prs)
    _header(slide, 11, "Memoria de Calculo — Equipamento Selecionado")
    _footer(slide)

    cy = HH + 0.16
    d  = res["selected_equipment_detail"]
    lw = 6.6
    rx = W - MX - 2.45

    # Info bar
    _rect(slide, MX, cy, lw, 0.50, fill=PAPER, line_color=LINE)
    _rect(slide, MX, cy, 0.06, 0.50, fill=BLUE)
    _txt(slide, f"{d['equipment_type']} — {d['supplier']} modelo {d['model']}",
         MX + 0.14, cy + 0.04, lw - 0.22, 0.24, size=11, bold=True, color=BLUE)
    _txt(slide, f"Operacao: {d['operations'][0]}  |  {d['reason']}",
         MX + 0.14, cy + 0.28, lw - 0.22, 0.18, size=8, color=MUTED)

    # Formula
    fy = cy + 0.58
    _rect(slide, MX, fy, lw, 0.40, fill=LGRAY, line_color=LINE)
    _txt(slide, "Formula:", MX + 0.1, fy + 0.03, 0.9, 0.16, size=7.5, bold=True, color=MUTED)
    _txt(slide, d["formula"], MX + 0.1, fy + 0.20, lw - 0.2, 0.17, size=8.5, bold=True, color=INK)

    # Inputs table
    ty = fy + 0.48
    headers = ["Variavel", "Valor", "Fonte / Premissa"]
    inputs = [
        ["Tempo-padrao por kit",    f"{d['standard_time_seconds_per_kit']:.0f} s",
         "Estimativa de engenharia"],
        ["Taxa nominal",            f"{d['nominal_rate_used']:.0f} kits/h",
         "3600 / tempo-padrao"],
        ["Horas uteis / semana",    f"{d['useful_hours_per_week']:.0f} h",
         "5 dias x 7 h uteis/dia"],
        ["Eficiencia",              f"{d['efficiency']*100:.0f}%",   "Premissa"],
        ["Confiabilidade",          f"{d['reliability']*100:.0f}%",  "Premissa"],
        ["Rendimento do processo",  f"{d['process_yield']*100:.0f}%","Premissa"],
        ["Demanda bruta semanal",   f"{d['demand_input_kits_per_week']} kits",
         "Meta boa / rendimento final"],
    ]
    _tbl(slide, headers, inputs, MX, ty, lw, FY - ty - 0.05,
         col_w=[2.7, 1.3, 2.6], row_size=8)

    # Right KPIs
    kh = (FY - cy - 0.14) / 2
    _kpi(slide, "Quantidade necessaria",
         str(d["required_quantity"]), "unidades",
         rx, cy, 2.1, kh, GREEN)
    _kpi(slide, "Utilizacao",
         f"{d['utilization']*100:.1f}%", "das 2 unidades",
         rx, cy + kh + 0.10, 2.1, kh - 0.05, BLUE)


def s12_layout(prs, renders):
    slide = _slide(prs)
    _header(slide, 12, "Layout Esquematico — 24 x 16 m = 384 m2")
    _footer(slide)
    png = renders.get("layout_render")
    if png and Path(png).exists():
        iy = HH + 0.08
        _img(slide, png, MX, iy, W - 2 * MX, FY - iy - 0.05)
    else:
        _txt(slide, "[Layout nao disponivel]", MX, HH + 1.0, W - 2 * MX, 0.5,
             size=14, color=MUTED, align=PP_ALIGN.CENTER)


def s13_mapo(prs, renders):
    slide = _slide(prs)
    _header(slide, 13, "Mapofluxograma — Fluxo sobre o Layout")
    _footer(slide)
    png = renders.get("mapofluxograma_render")
    if png and Path(png).exists():
        iy = HH + 0.08
        _img(slide, png, MX, iy, W - 2 * MX, FY - iy - 0.05)
    else:
        _txt(slide, "[Mapofluxograma nao disponivel]", MX, HH + 1.0, W - 2 * MX, 0.5,
             size=14, color=MUTED, align=PP_ALIGN.CENTER)


def s14_conclusoes(prs, proj, res):
    slide = _slide(prs)
    _header(slide, 14, "Conclusoes")
    _footer(slide)

    cy   = HH + 0.14
    conc = proj["conclusions"]
    lay  = res["layout"]
    bot  = res["bottleneck"]

    # Summary box
    _rect(slide, MX, cy, W - 2 * MX, 0.68, fill=PAPER, line_color=LINE)
    _txt(slide, conc["summary"],
         MX + 0.12, cy + 0.06, W - 2 * MX - 0.22, 0.56,
         size=9.5, color=INK, wrap=True)

    # KPIs
    ky = cy + 0.76
    kh = 0.80
    kw = (W - 2 * MX - 0.39) / 4
    kpis_c = [
        ("Area total",       "384 m2",                  "24 x 16 m",            BLUE),
        ("Ocupacao",         f"{lay['occupancy']*100:.1f}%", "67,8% ocupado",    GREEN),
        ("Gargalo",          "Router CNC",              "Maksiwa RTC.1313",      GOLD),
        ("Utilizacao garg.", f"{bot['utilization']*100:.1f}%", "2 unidades",     RED),
    ]
    for j, (lbl, val, unit, c) in enumerate(kpis_c):
        _kpi(slide, lbl, val, unit, MX + j * (kw + 0.13), ky, kw, kh, c)

    # Improvements
    iy2 = ky + kh + 0.12
    _txt(slide, "Oportunidades de melhoria:", MX, iy2, W - 2 * MX, 0.22,
         size=9, bold=True, color=INK)

    impr = conc["improvements"]
    hw   = (W - 2 * MX - 0.18) / 2
    _lines(slide, ["• " + x for x in impr[:3]],
           MX, iy2 + 0.24, hw, FY - iy2 - 0.28, size=8, color=INK, spacing_pt=2)
    _lines(slide, ["• " + x for x in impr[3:]],
           MX + hw + 0.18, iy2 + 0.24, hw, FY - iy2 - 0.28, size=8, color=INK, spacing_pt=2)


# ── main ─────────────────────────────────────────────────────────────────────

def main():
    print("Carregando dados...")
    proj = json.loads((DATA / "projeto.json").read_text(encoding="utf-8"))
    res  = json.loads((DATA / "resultados_calculo.json").read_text(encoding="utf-8"))

    print("Convertendo SVGs para PNG (Playwright)...")
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        try:
            renders = _render_pngs(tmp_path)
        except Exception as e:
            print(f"  Aviso: conversao SVG falhou ({e}). Slides de diagrama ficarão em branco.")
            renders = {}

        print("Montando apresentacao PPTX...")
        prs = Presentation()
        prs.slide_width  = Inches(W)
        prs.slide_height = Inches(H)

        s01_cover(prs, proj, res)
        s02_objetivos(prs, proj)
        s03_produto(prs, proj)
        s04_bom(prs, proj)
        s05_mercado(prs, proj)
        s06_meta(prs, proj, res)
        s07_tabela1(prs, proj)
        s08_processos(prs, proj)
        s09_fluxograma(prs, renders)
        s10_equipamentos(prs, res)
        s11_calculo(prs, res)
        s12_layout(prs, renders)
        s13_mapo(prs, renders)
        s14_conclusoes(prs, proj, res)

        OUT_DIR.mkdir(exist_ok=True)
        out = OUT_DIR / "apresentacao_tramontina_22399036.pptx"
        prs.save(str(out))
        print(f"\nPronto -> {out}")


if __name__ == "__main__":
    main()

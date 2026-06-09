# polish_presentation.py -- Polishes the group's PPTX with renders and data corrections.
#
# Input:  C:/Users/dvill/Downloads/Projeto de Fabrica - Kit Churrasco Tramontina (1).pptx
# Output: 01_apresentacao/apresentacao_final.pptx


import tempfile
from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT    = Path(__file__).parent.parent
RENDERS = ROOT / "06_dashboard" / "renders"
OUT_DIR = ROOT / "01_apresentacao"
SOURCE  = Path(r"C:\Users\dvill\Downloads") / "Projeto de Fábrica - Kit Churrasco Tramontina (1).pptx"

W = 10.0
H = 5.625


def _svg_to_png(svg: Path, png: Path, width=1800, height=1000):
    from playwright.sync_api import sync_playwright
    edge = Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe")
    kwargs = {"headless": True}
    if edge.exists():
        kwargs["executable_path"] = str(edge)
    with sync_playwright() as p:
        browser = p.chromium.launch(**kwargs)
        try:
            page = browser.new_page(viewport={"width": width, "height": height})
            page.goto(svg.as_uri(), wait_until="networkidle")
            page.screenshot(path=str(png))
        finally:
            browser.close()


def _split_png(src: Path, top: Path, bot: Path):
    """Split PNG: top covers 0–60%, bottom covers 40–100% (20% overlap)."""
    img = Image.open(str(src))
    w, h = img.size
    img.crop((0, 0,            w, int(h * 0.60))).save(str(top))
    img.crop((0, int(h * 0.40), w, h           )).save(str(bot))


def _render_pngs(tmp: Path) -> dict:
    pngs = {}
    print("  fluxograma_render.svg -> PNG ...")
    full = tmp / "fluxo_full.png"
    fluxo_svg = RENDERS / "fluxograma_render.svg"
    if not fluxo_svg.exists():
        raise FileNotFoundError(f"Missing render: {fluxo_svg}")
    _svg_to_png(fluxo_svg, full, width=1800, height=1000)
    pngs["fluxo_top"] = tmp / "fluxo_top.png"
    pngs["fluxo_bot"] = tmp / "fluxo_bot.png"
    _split_png(full, pngs["fluxo_top"], pngs["fluxo_bot"])
    for name in ("layout_render", "mapofluxograma_render"):
        print(f"  {name}.svg -> PNG ...")
        svg_path = RENDERS / f"{name}.svg"
        if not svg_path.exists():
            raise FileNotFoundError(f"Missing render: {svg_path}")
        png = tmp / f"{name}.png"
        _svg_to_png(svg_path, png)
        pngs[name] = png
    return pngs


def _add_img_slide(prs, idx: int, title: str, img_path: Path):
    """Insert a new blank slide with a full-width image at 0-based position idx."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank, added at end

    # Move from end to target idx
    lst = prs.slides._sldIdLst
    ref = lst[-1]
    lst.remove(ref)
    lst.insert(idx, ref)

    # Dark title bar at bottom
    bar_h = 0.40
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(H - bar_h), Inches(W), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x17, 0x21, 0x2b)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0), Inches(H - bar_h), Inches(W), Inches(bar_h))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    # Image fills area above bar
    slide.shapes.add_picture(
        str(img_path), Inches(0), Inches(0), Inches(W), Inches(H - bar_h - 0.02)
    )
    return slide


def _insert_fluxograma_slides(prs, pngs: dict):
    """Insert 2 fluxograma slides after index 10 ('PAIS E FILHOS')."""
    _add_img_slide(prs, 11, "Fluxograma do Processo (1/2)", pngs["fluxo_top"])
    _add_img_slide(prs, 12, "Fluxograma do Processo (2/2)", pngs["fluxo_bot"])


def _fill_existing_slide(slide, img_path: Path):
    """Add image below the lowest existing shape on the slide."""
    lowest_emu = 0
    for shape in slide.shapes:
        bottom = shape.top + shape.height
        if bottom > lowest_emu:
            lowest_emu = bottom
    top_in = max(lowest_emu / 914400 + 0.08, 0.75)
    slide.shapes.add_picture(
        str(img_path),
        Inches(0.1), Inches(top_in),
        Inches(W - 0.2), Inches(H - top_in - 0.08)
    )


# Longer/more-specific strings first to avoid partial-match shadowing
TEXT_FIXES = [
    # Meta — production target
    ("5.000 kits semanais",        "1.000 kits semanais"),
    ("5.000 kits por semana",      "1.000 kits por semana"),
    ("2.000 kits por semana",      "1.000 kits por semana"),
    ("Calcular a carga de máquinas e layout para atender à meta de 5.000 kits semanais.",
     "Calcular a carga de máquinas e layout para atender à meta de 1.000 kits semanais."),
    ("meta de 2.000 unidades/semana (400/dia + rejeitos)",
     "meta de 1.000 kits/semana (200/dia + rejeitos)"),
    ("Cálculo da quantidade de prensas para meta de 2.000 unidades/semana (400/dia + rejeitos).",
     "Cálculo da quantidade de lasers para meta de 1.000 kits/semana (200/dia + rejeitos)."),
    ("2.000",                      "1.000"),
    ("400 kits / dia",             "200 kits / dia"),
    ("400 kits/dia",               "200 kits/dia"),
    ("Capacidade Diária: 400 kits.", "Capacidade Diária: 200 kits."),
    ("400 / 8 = 50",               "200 / 8 = 25"),
    ("2000 kits / 5 dias = 400 kits / dia",
     "1000 kits / 5 dias = 200 kits / dia"),
    ("50 / 0,85 x 0,97 = 60,7 kits / hora",
     "25 / 0,85 x 0,97 = 30,3 kits / hora"),
    ("50 kits / hora",             "30 kits / hora"),
    ("60,7 kits / hora",           "30,3 kits / hora"),
    # Layout
    ("800m²",                      "384 m²"),
    ("800 m²",                     "384 m²"),
    ("40m x 20m",                  "24 × 16 m"),
    ("40 × 20",                    "24 × 16"),
    ("Planta dimensionada em 800m² (40m x 20m) focada em fluxo linear para minimizar movimentação de materiais.",
     "Planta dimensionada em 384 m² (24 × 16 m) com arranjo misto: setores funcionais para metal/madeira e linha para montagem/embalagem."),
    # Material
    ("350 x 220 x 18 mm",         "340 x 190 x 15 mm"),
    ("Madeira Certificada (Eucalipto/Pinus)", "Madeira Maçaranduba"),
    ("Eucalipto/Pinus",            "Maçaranduba"),
    # Equipment — specific phrases before generic
    ("A Prensa Excêntrica foi dimensionada para as operações de corte e estampagem dos componentes metálicos. Considerando a demanda de",
     "O Laser Fibra CNC foi dimensionado para o corte dos blanks de faca e garfo em chapa AISI 420. Considerando a demanda de"),
    ("1 Prensa Excêntrica é suficiente para a demanda, operando com uma folga de capacidade de 70%.",
     "1 Laser Fibra CNC é suficiente para a demanda, operando com utilização de 49%."),
    ("Harlo do Brasil\n(Guarulhos-SP)\nRHTC / ProfiPress", "Madetech (SP)"),
    ("Harlo do Brasil",            "Madetech (SP)"),
    ("60 golpes/min",              "20.000 mm/min"),
    ("60 tf",                      "1.500 W"),
    ("Prensa\nExcêntrica",         "Laser Fibra CNC"),
    ("Prensa Excêntrica",          "Laser Fibra CNC"),
    ("prensas para",               "lasers para"),
    ("quantidade de prensas",      "quantidade de lasers"),
    ("1 prensa",                   "1 laser"),
    ("da prensa",                  "do laser"),
    # Calculation numbers
    ("N = 60,7 /3600 = 0,017 ≈ 1 prensa",  "N = 30,3 / 80 = 0,38 ≈ 1 laser"),
    ("N = 412,37 / 1360 = 0,3 -> 1 prensa", "N = 206,2 / 437 = 0,47 → 1 laser"),
    ("Db = 400 / 0,97 = 412,37 kits/dia",   "Db = 200 / 0,97 = 206,2 kits/dia"),
    ("Capacidade = 367,2 / 0,27 = 1360 kits / dia",
     "Capacidade = 80 × 7 × 0,85 × 0,92 = 437 kits/dia"),
    ("Capacidade da prensa:  60 golpes / min = 3600 ciclos / hora",
     "Capacidade do laser: 20.000 mm/min → ~80 kits/hora"),
]


def _replace_in_para(para, fixes):
    """Two-pass replacement: run-level first, then paragraph-level for split runs."""
    for run in para.runs:
        for old, new in fixes:
            if old in run.text:
                run.text = run.text.replace(old, new)
    # Paragraph-level pass catches strings split across runs
    full = "".join(r.text for r in para.runs)
    new_full = full
    for old, new in fixes:
        if old in new_full:
            new_full = new_full.replace(old, new)
    if new_full != full and para.runs:
        para.runs[0].text = new_full
        for r in para.runs[1:]:
            r.text = ""


def _fix_texts(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    _replace_in_para(para, TEXT_FIXES)
            if hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            _replace_in_para(para, TEXT_FIXES)


def _fill_diagram_slides(prs, pngs: dict):
    """Find diagram placeholder slides by title keyword and add renders."""
    targets = {
        "ESQUEMÁTICO": pngs["layout_render"],
        "MAPOFLUXOGRAMA": pngs["mapofluxograma_render"],
    }
    filled: set = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.upper()
            for keyword, img_path in targets.items():
                if keyword in text and keyword not in filled:
                    print(f"  Filling: {shape.text_frame.text.strip()[:50]}")
                    _fill_existing_slide(slide, img_path)
                    filled.add(keyword)
                    break


def main():
    print("Rendering SVGs to PNG...")
    with tempfile.TemporaryDirectory() as tmp:
        pngs = _render_pngs(Path(tmp))

        print("Loading source PPTX...")
        prs = Presentation(str(SOURCE))
        print(f"  Slides before: {len(prs.slides)}")

        print("Inserting fluxograma slides...")
        _insert_fluxograma_slides(prs, pngs)
        print(f"  Slides after: {len(prs.slides)}")

        print("Filling diagram slides...")
        _fill_diagram_slides(prs, pngs)

        print("Applying text corrections...")
        _fix_texts(prs)

        OUT_DIR.mkdir(exist_ok=True)
        out = OUT_DIR / "apresentacao_final.pptx"
        prs.save(str(out))
        print(f"Saved -> {out}")


if __name__ == "__main__":
    main()
